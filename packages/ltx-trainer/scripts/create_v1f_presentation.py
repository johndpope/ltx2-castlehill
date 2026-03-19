"""Generate VFM v1f PowerPoint presentation with illustrations."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Color palette ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)      # dark navy
BG_CARD = RGBColor(0x16, 0x20, 0x3A)      # card background
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)  # bright cyan
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA) # lavender
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)  # mint green
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24) # amber
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)    # coral
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16, color=WHITE, line_spacing=1.2, font_name="Calibri"):
    """lines: list of (text, color, bold, font_size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, c, b, fs = item, color, False, font_size
        elif len(item) == 2:
            text, c = item
            b, fs = False, font_size
        elif len(item) == 3:
            text, c, b = item
            fs = font_size
        else:
            text, c, b, fs = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT_BLUE, width=Pt(2)):
    """Add a connector arrow."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2,
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def draw_box_with_label(slide, left, top, width, height, label, fill=BG_CARD, border=ACCENT_BLUE, font_size=14, text_color=WHITE):
    """Draw a rounded rect with centered label."""
    shape = add_shape(slide, left, top, width, height, fill_color=fill, line_color=border, line_width=Pt(2))
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(4)
    return shape


def draw_down_arrow(slide, cx, y_start, length, color=ACCENT_BLUE):
    """Draw a downward arrow at center x."""
    # Shaft
    shaft = add_shape(slide, cx - Inches(0.03), y_start, Inches(0.06), length,
                      fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    # Arrowhead
    tri = add_shape(slide, cx - Inches(0.15), y_start + length - Inches(0.05),
                    Inches(0.3), Inches(0.2),
                    fill_color=color, shape_type=MSO_SHAPE.ISOSCELES_TRIANGLE)
    return shaft, tri


def draw_right_arrow_shape(slide, x, y, length, color=ACCENT_BLUE):
    """Draw horizontal arrow."""
    shaft = add_shape(slide, x, y - Inches(0.03), length, Inches(0.06),
                      fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    tri = add_shape(slide, x + length - Inches(0.05), y - Inches(0.12),
                    Inches(0.2), Inches(0.24),
                    fill_color=color, shape_type=MSO_SHAPE.RIGHT_TRIANGLE)
    return shaft


# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Decorative sphere circles
for i, (cx, cy, r, clr) in enumerate([
    (Inches(10.5), Inches(1.5), Inches(2.5), ACCENT_BLUE),
    (Inches(11.2), Inches(2.2), Inches(1.8), ACCENT_PURPLE),
    (Inches(10.0), Inches(2.8), Inches(1.2), ACCENT_GREEN),
]):
    circle = add_shape(slide, cx - r//2, cy - r//2, r, r,
                       line_color=clr, line_width=Pt(1.5),
                       shape_type=MSO_SHAPE.OVAL)
    circle.fill.background()

add_text(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.6),
         "VFM v1f", font_size=20, color=ACCENT_BLUE, bold=True)

add_text(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(1.5),
         "Spherical Cauchy Noise Adapter", font_size=52, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(3.3), Inches(10), Inches(1.0),
         "One-Step Video Generation via Structured Noise on S\u00b9\u00b2\u2077",
         font_size=28, color=GRAY)

# Key stats boxes
stats = [
    ("8.7\u201321.5x", "Faster than\n8-step LTX-2", ACCENT_BLUE),
    ("1 Step", "Forward pass\n(was 8 steps)", ACCENT_GREEN),
    ("S\u00b9\u00b2\u2077", "Noise lives on\nhypersphere", ACCENT_PURPLE),
    ("19.1M", "Adapter\nparameters", ACCENT_ORANGE),
]
for i, (val, desc, clr) in enumerate(stats):
    x = Inches(0.8 + i * 2.8)
    y = Inches(4.8)
    box = add_shape(slide, x, y, Inches(2.4), Inches(1.8),
                    fill_color=BG_CARD, line_color=clr, line_width=Pt(2))
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.0), Inches(0.6),
             val, font_size=32, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.0), Inches(0.8),
             desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM — Why 1-step generation?
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
         "The Problem: 8 Steps is Too Slow", font_size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Standard LTX-2 requires 8 Euler denoising steps \u2014 each step runs the full 48-layer, 19B DiT",
         font_size=18, color=GRAY)

# Standard pipeline (slow)
y_base = Inches(2.2)
add_text(slide, Inches(0.8), y_base - Inches(0.4), Inches(4), Inches(0.4),
         "Standard LTX-2  (8 steps)", font_size=20, color=ACCENT_RED, bold=True)

box_w = Inches(1.2)
box_h = Inches(0.7)
for step in range(8):
    x = Inches(0.8 + step * 1.4)
    box = draw_box_with_label(slide, x, y_base, box_w, box_h,
                              f"Step {step+1}\n48 layers", fill=BG_CARD, border=ACCENT_RED, font_size=11)
    if step < 7:
        draw_right_arrow_shape(slide, x + box_w, y_base + box_h/2, Inches(0.2), color=ACCENT_RED)

add_text(slide, Inches(0.8), y_base + Inches(0.8), Inches(11), Inches(0.4),
         "z ~ N(0, I)  \u2192  [step 1] \u2192 [step 2] \u2192 ... \u2192 [step 8] \u2192  video     ~33\u201376 seconds",
         font_size=16, color=LIGHT_GRAY, font_name="Consolas")

# VFM pipeline (fast)
y_vfm = Inches(4.0)
add_text(slide, Inches(0.8), y_vfm - Inches(0.4), Inches(4), Inches(0.4),
         "VFM v1f  (1 step)", font_size=20, color=ACCENT_GREEN, bold=True)

# Adapter box
draw_box_with_label(slide, Inches(0.8), y_vfm, Inches(2.5), Inches(1.2),
                    "Noise Adapter q\u03c6\n\nText \u2192 Structured Noise\non S\u00b9\u00b2\u2077",
                    fill=BG_CARD, border=ACCENT_PURPLE, font_size=13)

draw_right_arrow_shape(slide, Inches(3.3), y_vfm + Inches(0.6), Inches(0.5), color=ACCENT_GREEN)

# Single DiT step
draw_box_with_label(slide, Inches(3.9), y_vfm, Inches(2.5), Inches(1.2),
                    "1 Step\n48-layer DiT\n(LoRA fine-tuned)",
                    fill=BG_CARD, border=ACCENT_GREEN, font_size=13)

draw_right_arrow_shape(slide, Inches(6.4), y_vfm + Inches(0.6), Inches(0.5), color=ACCENT_GREEN)

# Output
draw_box_with_label(slide, Inches(7.0), y_vfm, Inches(2.0), Inches(1.2),
                    "Video\n768\u00d7448\n4 frames",
                    fill=BG_CARD, border=ACCENT_GREEN, font_size=13)

add_text(slide, Inches(0.8), y_vfm + Inches(1.3), Inches(11), Inches(0.4),
         "z ~ q\u03c6(z|text)  \u2192  [1 step \u00d7 48 layers]  \u2192  video     ~3.5\u20138 seconds",
         font_size=16, color=LIGHT_GRAY, font_name="Consolas")

# Key insight box
insight_box = add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2),
                        fill_color=RGBColor(0x1A, 0x2A, 0x4A), line_color=ACCENT_BLUE, line_width=Pt(2))
add_multiline(slide, Inches(1.2), Inches(5.9), Inches(10.5), Inches(1.0), [
    ("Key Insight: The noise IS the information", ACCENT_BLUE, True, 20),
    ("Instead of random noise that requires 8 steps to refine, learn noise that already encodes the target video's structure. The 48-layer DiT then acts as a single-step decoder.", LIGHT_GRAY, False, 16),
])


# ════════════════════════════════════════════════════════════════
# SLIDE 3: VFM VERSION EVOLUTION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
         "VFM Evolution: v1a \u2192 v1f", font_size=40, color=WHITE, bold=True)

versions = [
    ("v1a", "Vanilla MLP", "Pooled text \u2192 same noise for all tokens", GRAY, "Baseline"),
    ("v1b", "Transformer", "Cross-attn + positions \u2192 per-token noise", GRAY, "+19.1M params"),
    ("v1c", "Diversity", "Temporal + spatial diversity regularization", GRAY, "Anti-collapse"),
    ("v1d", "Trajectory\nDistill", "Per-token \u03c3 + teacher ODE trajectory matching", ACCENT_BLUE, "+SigmaHead"),
    ("v1e", "Content\nRouter", "Complexity-adaptive \u03c3 from GT latent features", GRAY, "+Router 1.6M"),
    ("v1f", "Spherical\nCauchy", "Noise on S\u00b9\u00b2\u2077: direction (\u03bc\u0302) + magnitude + \u03ba", ACCENT_GREEN, "This work"),
]

for i, (ver, name, desc, clr, tag) in enumerate(versions):
    y = Inches(1.5 + i * 0.9)
    # Version badge
    badge_clr = ACCENT_GREEN if ver == "v1f" else (ACCENT_BLUE if ver == "v1d" else BG_CARD)
    badge = add_shape(slide, Inches(0.8), y, Inches(0.9), Inches(0.6),
                      fill_color=badge_clr, line_color=clr if clr != GRAY else RGBColor(0x47, 0x55, 0x69))
    p = badge.text_frame.paragraphs[0]
    p.text = ver
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Name
    add_text(slide, Inches(1.9), y, Inches(1.8), Inches(0.6),
             name, font_size=15, color=WHITE if clr != GRAY else LIGHT_GRAY, bold=True)

    # Description
    add_text(slide, Inches(3.8), y, Inches(6.5), Inches(0.6),
             desc, font_size=14, color=LIGHT_GRAY if clr != GRAY else GRAY)

    # Tag
    add_text(slide, Inches(10.5), y, Inches(2.0), Inches(0.6),
             tag, font_size=12, color=clr if clr != GRAY else RGBColor(0x64, 0x74, 0x8B),
             alignment=PP_ALIGN.RIGHT)

    # Arrow connecting versions
    if i < len(versions) - 1:
        draw_down_arrow(slide, Inches(1.25), y + Inches(0.6), Inches(0.25),
                        color=RGBColor(0x33, 0x44, 0x55))

# v1f highlight box
highlight = add_shape(slide, Inches(0.6), Inches(5.8) + Inches(0.15), Inches(11.8), Inches(0.8),
                      line_color=ACCENT_GREEN, line_width=Pt(2))

# Bottom note
add_text(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
         "v1f branches from v1d (not v1e) \u2014 same adapter architecture, reinterpreted outputs. Weights transfer directly.",
         font_size=14, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: GAUSSIAN vs SPHERICAL CAUCHY
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "Gaussian vs Spherical Cauchy Noise", font_size=40, color=WHITE, bold=True)

# Left: Gaussian
left_x = Inches(0.8)
add_text(slide, left_x, Inches(1.3), Inches(5.5), Inches(0.5),
         "v1d: Gaussian Reparameterization", font_size=22, color=ACCENT_RED, bold=True)

gauss_box = add_shape(slide, left_x, Inches(1.9), Inches(5.5), Inches(4.8),
                      fill_color=BG_CARD, line_color=RGBColor(0x47, 0x55, 0x69), line_width=Pt(1))

add_multiline(slide, left_x + Inches(0.3), Inches(2.1), Inches(4.9), Inches(4.5), [
    ("z = \u03bc + \u03c3 \u00b7 \u03b5,  \u03b5 ~ N(0, I)", ACCENT_RED, True, 20),
    ("", WHITE, False, 8),
    ("Properties:", LIGHT_GRAY, True, 16),
    ("\u2022 Noise lives in unbounded R\u00b9\u00b2\u2078", GRAY, False, 14),
    ("\u2022 Direction and magnitude are entangled", GRAY, False, 14),
    ("\u2022 Gaussian tails decay exponentially", GRAY, False, 14),
    ("\u2022 KL = \u00bd[\u03bc\u00b2 + \u03c3\u00b2 - log(\u03c3\u00b2) - 1] per dim", GRAY, False, 14),
    ("", WHITE, False, 8),
    ("Problems:", ACCENT_RED, True, 16),
    ("\u2022 All 128 dims treated equally", GRAY, False, 14),
    ("\u2022 No geometric structure to noise", GRAY, False, 14),
    ("\u2022 Hard to measure diversity meaningfully", GRAY, False, 14),
    ("\u2022 \u03c3 collapse \u2192 delta function (mode collapse)", GRAY, False, 14),
])

# Right: Spherical Cauchy
right_x = Inches(7.0)
add_text(slide, right_x, Inches(1.3), Inches(5.5), Inches(0.5),
         "v1f: Spherical Cauchy on S\u00b9\u00b2\u2077", font_size=22, color=ACCENT_GREEN, bold=True)

sph_box = add_shape(slide, right_x, Inches(1.9), Inches(5.5), Inches(4.8),
                    fill_color=BG_CARD, line_color=ACCENT_GREEN, line_width=Pt(1))

add_multiline(slide, right_x + Inches(0.3), Inches(2.1), Inches(4.9), Inches(4.5), [
    ("z = \u2016\u03bc\u2016 \u00b7 SpCauchy(\u03bc\u0302, \u03ba)", ACCENT_GREEN, True, 20),
    ("", WHITE, False, 8),
    ("Properties:", LIGHT_GRAY, True, 16),
    ("\u2022 Direction lives on unit hypersphere S\u00b9\u00b2\u2077", GRAY, False, 14),
    ("\u2022 Direction and magnitude are separated", GRAY, False, 14),
    ("\u2022 Heavy tails \u2192 broader exploration", GRAY, False, 14),
    ("\u2022 KL = (D-1)/2 \u00b7 log(1 + 1/\u03ba) \u2014 clean!", GRAY, False, 14),
    ("", WHITE, False, 8),
    ("Advantages:", ACCENT_GREEN, True, 16),
    ("\u2022 \u03bc\u0302 = what kind of noise (direction)", GRAY, False, 14),
    ("\u2022 \u2016\u03bc\u2016 = how much noise (magnitude)", GRAY, False, 14),
    ("\u2022 \u03ba = how confident (concentration/token)", GRAY, False, 14),
    ("\u2022 Geodesic distance = meaningful diversity", GRAY, False, 14),
])

# VS circle
vs_circle = add_shape(slide, Inches(5.9), Inches(3.8), Inches(0.8), Inches(0.8),
                      fill_color=BG_DARK, line_color=ACCENT_PURPLE, line_width=Pt(2),
                      shape_type=MSO_SHAPE.OVAL)
p = vs_circle.text_frame.paragraphs[0]
p.text = "vs"
p.font.size = Pt(18)
p.font.color.rgb = ACCENT_PURPLE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════
# SLIDE 5: DIRECTION-MAGNITUDE DECOMPOSITION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "Direction-Magnitude Decomposition", font_size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Same adapter architecture (NoiseAdapterV1b) \u2014 outputs reinterpreted geometrically. Weights transfer from v1d.",
         font_size=16, color=GRAY)

# Adapter output
adapter_box = draw_box_with_label(slide, Inches(0.8), Inches(2.0), Inches(3.0), Inches(1.0),
                                  "Adapter Output\n\u03bc [B, 1344, 128]  log_\u03c3 [B, 1344, 128]",
                                  fill=BG_CARD, border=ACCENT_PURPLE, font_size=13)

# Three branches
branches = [
    ("\u03bc\u0302 = normalize(\u03bc)", "Direction on S\u00b9\u00b2\u2077\n\"What kind of noise\"", ACCENT_BLUE, Inches(0.5)),
    ("r = \u2016\u03bc\u2016", "Magnitude (scalar/token)\n\"How much noise\"", ACCENT_GREEN, Inches(4.5)),
    ("\u03ba = exp(mean(log_\u03c3))", "Concentration (scalar/token)\n\"How confident\"", ACCENT_ORANGE, Inches(8.5)),
]

for label, desc, clr, x in branches:
    # Arrow down from adapter
    draw_down_arrow(slide, x + Inches(1.5), Inches(3.0), Inches(0.5), color=clr)

    # Branch box
    branch_box = add_shape(slide, x, Inches(3.7), Inches(3.0), Inches(1.5),
                           fill_color=BG_CARD, line_color=clr, line_width=Pt(2))

    add_text(slide, x + Inches(0.2), Inches(3.8), Inches(2.6), Inches(0.5),
             label, font_size=18, color=clr, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")
    add_text(slide, x + Inches(0.2), Inches(4.35), Inches(2.6), Inches(0.7),
             desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Final combination
for x in [Inches(2.0), Inches(6.0), Inches(10.0)]:
    draw_down_arrow(slide, x, Inches(5.2), Inches(0.4), color=ACCENT_PURPLE)

combo_box = add_shape(slide, Inches(1.5), Inches(5.8), Inches(9.5), Inches(1.2),
                      fill_color=RGBColor(0x1A, 0x2A, 0x4A), line_color=ACCENT_PURPLE, line_width=Pt(2))

add_multiline(slide, Inches(1.8), Inches(5.9), Inches(9.0), Inches(1.0), [
    ("z_dir ~ SphericalCauchy(\u03bc\u0302, \u03ba)     \u2192     z = r \u00b7 z_dir", ACCENT_PURPLE, True, 22),
    ("Sample direction on hypersphere, then scale by learned magnitude per token", GRAY, False, 14),
])


# ════════════════════════════════════════════════════════════════
# SLIDE 6: SPHERICAL CAUCHY DISTRIBUTION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "Spherical Cauchy Distribution on S\u00b9\u00b2\u2077", font_size=40, color=WHITE, bold=True)

# Sampling algorithm
algo_box = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
                     fill_color=BG_CARD, line_color=ACCENT_BLUE, line_width=Pt(1))

add_text(slide, Inches(1.0), Inches(1.6), Inches(5.0), Inches(0.4),
         "Sampling Algorithm", font_size=20, color=ACCENT_BLUE, bold=True)

add_multiline(slide, Inches(1.0), Inches(2.1), Inches(5.0), Inches(4.2), [
    ("1. Inverse CDF for Cauchy:", LIGHT_GRAY, True, 15),
    ("   t = tan(\u03c0(u - 0.5)),  u ~ Uniform(0,1)", WHITE, False, 14),
    ("", WHITE, False, 6),
    ("2. Scale by concentration:", LIGHT_GRAY, True, 15),
    ("   r = \u221a\u03ba \u00b7 t,  clamp to [-10, 10]", WHITE, False, 14),
    ("", WHITE, False, 6),
    ("3. Random tangent vector via Gram-Schmidt:", LIGHT_GRAY, True, 15),
    ("   tangent = noise - (noise \u00b7 \u03bc\u0302)\u03bc\u0302", WHITE, False, 14),
    ("   tangent = normalize(tangent)", WHITE, False, 14),
    ("", WHITE, False, 6),
    ("4. Exponential map:", LIGHT_GRAY, True, 15),
    ("   sample = cos(r)\u00b7\u03bc\u0302 + sin(r)\u00b7tangent", WHITE, False, 14),
    ("   sample = normalize(sample)", WHITE, False, 14),
], font_name="Consolas")

# Kappa behavior
kappa_box = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.2),
                      fill_color=BG_CARD, line_color=ACCENT_ORANGE, line_width=Pt(1))

add_text(slide, Inches(7.0), Inches(1.6), Inches(5.0), Inches(0.4),
         "\u03ba Controls Concentration", font_size=20, color=ACCENT_ORANGE, bold=True)

kappa_vals = [
    ("\u03ba \u2192 0", "Peaked at \u03bc\u0302 (delta-like)", "KL \u2192 \u221e", ACCENT_RED),
    ("\u03ba = 1.0", "Moderate spread", "KL = 44.2", ACCENT_BLUE),
    ("\u03ba = 2.0", "Target concentration", "KL = 25.2", ACCENT_GREEN),
    ("\u03ba \u2192 \u221e", "Uniform on sphere", "KL \u2192 0", GRAY),
]

for i, (kval, desc, kl, clr) in enumerate(kappa_vals):
    y = Inches(2.15 + i * 0.35)
    add_text(slide, Inches(7.0), y, Inches(1.5), Inches(0.35),
             kval, font_size=13, color=clr, bold=True, font_name="Consolas")
    add_text(slide, Inches(8.5), y, Inches(2.0), Inches(0.35),
             desc, font_size=13, color=LIGHT_GRAY)
    add_text(slide, Inches(10.8), y, Inches(1.5), Inches(0.35),
             kl, font_size=13, color=clr, font_name="Consolas")

# KL formula
kl_box = add_shape(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(2.5),
                   fill_color=BG_CARD, line_color=ACCENT_GREEN, line_width=Pt(1))

add_text(slide, Inches(7.0), Inches(4.1), Inches(5.0), Inches(0.4),
         "Spherical KL Divergence", font_size=20, color=ACCENT_GREEN, bold=True)

add_multiline(slide, Inches(7.0), Inches(4.6), Inches(5.0), Inches(1.8), [
    ("KL(SpCauchy || Uniform)", LIGHT_GRAY, True, 16),
    ("  = (D-1)/2 \u00b7 log(1 + 1/\u03ba)", WHITE, True, 20),
    ("", WHITE, False, 6),
    ("D = 128 (latent dimension)", GRAY, False, 14),
    ("Always \u2265 0, monotonically decreasing in \u03ba", GRAY, False, 14),
    ("Higher \u03ba = more spread = lower KL", GRAY, False, 14),
], font_name="Consolas")


# ════════════════════════════════════════════════════════════════
# SLIDE 7: FULL ARCHITECTURE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.6),
         "v1f Full Architecture", font_size=36, color=WHITE, bold=True)

# Pipeline boxes
pipeline = [
    ("Text Encoder\n(Gemma, frozen)", Inches(0.5), Inches(1.3), ACCENT_BLUE, Inches(2.2)),
    ("Noise Adapter q\u03c6\n(4-layer Transformer)\n19.1M params", Inches(0.5), Inches(2.7), ACCENT_PURPLE, Inches(2.2)),
    ("Reparameterize\n\u03bc\u0302 = norm(\u03bc), r = \u2016\u03bc\u2016\n\u03ba = exp(mean(log_\u03c3))\nz = r \u00b7 SpCauchy(\u03bc\u0302, \u03ba)", Inches(0.5), Inches(4.5), ACCENT_GREEN, Inches(2.2)),
]

for label, x, y, clr, w in pipeline:
    draw_box_with_label(slide, x, y, w, Inches(1.0), label, fill=BG_CARD, border=clr, font_size=11)

# Arrows between pipeline
draw_down_arrow(slide, Inches(1.6), Inches(2.3), Inches(0.3), color=ACCENT_BLUE)
draw_down_arrow(slide, Inches(1.6), Inches(3.7), Inches(0.6), color=ACCENT_PURPLE)

# SigmaHead branch
draw_right_arrow_shape(slide, Inches(2.7), Inches(3.2), Inches(0.8), color=ACCENT_ORANGE)
draw_box_with_label(slide, Inches(3.6), Inches(2.8), Inches(2.0), Inches(0.8),
                    "SigmaHead\n128\u2192256\u2192256\u21921\n99K params", fill=BG_CARD, border=ACCENT_ORANGE, font_size=11)

# Per-token sigma output
draw_down_arrow(slide, Inches(4.6), Inches(3.6), Inches(0.3), color=ACCENT_ORANGE)
add_text(slide, Inches(3.8), Inches(4.0), Inches(2.0), Inches(0.3),
         "\u03c3_i \u2208 [0.05, 0.95]", font_size=12, color=ACCENT_ORANGE, font_name="Consolas")

# Interpolation
draw_box_with_label(slide, Inches(3.6), Inches(4.5), Inches(3.0), Inches(1.0),
                    "Noisy Interpolation\nx_t[i] = (1-\u03c3_i)\u00b7x\u2080[i] + \u03c3_i\u00b7z[i]\nPer-token noise levels",
                    fill=BG_CARD, border=ACCENT_ORANGE, font_size=11)

# GT input
draw_box_with_label(slide, Inches(3.6), Inches(5.8), Inches(1.5), Inches(0.7),
                    "GT x\u2080\n(latents)", fill=BG_CARD, border=GRAY, font_size=11)
draw_down_arrow(slide, Inches(4.35), Inches(5.5), Inches(0.2), color=GRAY)

# Arrow from reparam to interpolation
draw_right_arrow_shape(slide, Inches(2.7), Inches(5.0), Inches(0.8), color=ACCENT_GREEN)

# 48-layer DiT
draw_right_arrow_shape(slide, Inches(6.6), Inches(5.0), Inches(0.5), color=WHITE)
draw_box_with_label(slide, Inches(7.2), Inches(4.3), Inches(2.5), Inches(1.4),
                    "48-Layer LTX-2 DiT\n(LoRA fine-tuned)\nFlow Map f\u03b8\n~233M trainable",
                    fill=BG_CARD, border=WHITE, font_size=12)

# Output
draw_right_arrow_shape(slide, Inches(9.7), Inches(5.0), Inches(0.5), color=WHITE)
draw_box_with_label(slide, Inches(10.3), Inches(4.3), Inches(2.3), Inches(1.4),
                    "x\u0302\u2080 = z - v\n(1-step decode)\n\u2193\nVAE Decoder\n\u2193\nVideo",
                    fill=BG_CARD, border=ACCENT_GREEN, font_size=11)

# Loss components at bottom
add_text(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.3),
         "Loss Components:", font_size=16, color=WHITE, bold=True)

losses = [
    ("L_mf", "Flow matching MSE", ACCENT_BLUE),
    ("L_kl", "Spherical KL", ACCENT_GREEN),
    ("L_div", "Diversity", ACCENT_PURPLE),
    ("L_mag", "\u2016\u03bc\u2016 regularization", ACCENT_ORANGE),
    ("L_\u03ba", "Kappa pull+entropy", ACCENT_RED),
    ("L_\u03c3", "Sigma pull+entropy", LIGHT_GRAY),
]
for i, (name, desc, clr) in enumerate(losses):
    x = Inches(0.5 + i * 2.1)
    add_text(slide, x, Inches(6.35), Inches(2.0), Inches(0.3),
             f"{name}: {desc}", font_size=11, color=clr)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: PER-TOKEN CONCENTRATION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "Per-Token Concentration: \u03ba per Spatial Location", font_size=36, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Each of 1344 tokens gets its own \u03ba (concentration) and \u03c3 (noise level) \u2014 heterogeneous denoising",
         font_size=16, color=GRAY)

# Frame grid illustration
frame_labels = ["Frame 0", "Frame 1", "Frame 2", "Frame 3"]
for f_idx, f_label in enumerate(frame_labels):
    x_base = Inches(0.8 + f_idx * 3.0)
    add_text(slide, x_base, Inches(1.8), Inches(2.6), Inches(0.4),
             f_label, font_size=16, color=LIGHT_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    # Grid of tokens (simplified 4x6)
    for row in range(4):
        for col in range(6):
            tx = x_base + Inches(col * 0.43)
            ty = Inches(2.3 + row * 0.43)

            # Simulate different kappa regions
            # Edges and complex regions get high kappa (peaked), flat regions get low kappa (broad)
            is_edge = (row == 0 or row == 3 or col == 0 or col == 5)
            is_center = (row in [1, 2] and col in [2, 3])

            if is_edge:
                clr = ACCENT_RED  # High kappa (concentrated)
            elif is_center:
                clr = ACCENT_BLUE  # Low kappa (broad)
            else:
                clr = ACCENT_PURPLE  # Medium kappa

            cell = add_shape(slide, tx, ty, Inches(0.4), Inches(0.4),
                            fill_color=clr, shape_type=MSO_SHAPE.RECTANGLE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = clr

# Legend
legend_y = Inches(4.5)
legend_items = [
    (ACCENT_RED, "High \u03ba (concentrated)", "Edges, textures \u2014 precise noise direction matters"),
    (ACCENT_PURPLE, "Medium \u03ba", "Moderate regions \u2014 balanced exploration"),
    (ACCENT_BLUE, "Low \u03ba (broad)", "Flat regions (sky, walls) \u2014 any direction works"),
]

for i, (clr, label, desc) in enumerate(legend_items):
    y = legend_y + Inches(i * 0.6)
    cell = add_shape(slide, Inches(1.0), y, Inches(0.3), Inches(0.3),
                     fill_color=clr, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(1.5), y, Inches(2.5), Inches(0.3),
             label, font_size=15, color=clr, bold=True)
    add_text(slide, Inches(4.2), y, Inches(7.0), Inches(0.3),
             desc, font_size=14, color=GRAY)

# Bottom insight
insight = add_shape(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
                    fill_color=RGBColor(0x1A, 0x2A, 0x4A), line_color=ACCENT_GREEN, line_width=Pt(2))
add_multiline(slide, Inches(1.2), Inches(6.35), Inches(10.5), Inches(0.7), [
    ("Key: The adapter learns WHERE to be precise vs WHERE to explore", ACCENT_GREEN, True, 18),
    ("Faces/edges need exact noise direction (\u03ba\u2191), sky/background can use any direction (\u03ba\u2193)", GRAY, False, 14),
])


# ════════════════════════════════════════════════════════════════
# SLIDE 9: LOSS FUNCTION DEEP DIVE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.6),
         "v1f Loss Function", font_size=36, color=WHITE, bold=True)

# Main loss equation
eq_box = add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
                   fill_color=BG_CARD, line_color=ACCENT_BLUE, line_width=Pt(2))
add_text(slide, Inches(1.0), Inches(1.2), Inches(11.0), Inches(0.5),
         "L = L_mf + w_kl\u00b7L_kl + L_div + w_mag\u00b7L_mag + w_\u03ba\u00b7L_\u03ba + w_\u03c3\u00b7L_\u03c3",
         font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Loss component cards
loss_cards = [
    ("L_mf", "Flow Matching MSE", "MSE(v_pred, v_target)\nv = z - x\u2080", "3.0", ACCENT_BLUE,
     "Core loss: predicted velocity\nmatches target velocity"),
    ("L_kl", "Spherical KL", "(D-1)/2 \u00b7 log(1+1/\u03ba)\nw = 3.0, free_bits = 0.05", "3.0", ACCENT_GREEN,
     "Regularizes noise distribution\nPrevents adapter collapse"),
    ("L_mag", "Magnitude Reg", "(\u2016\u03bc\u2016 - 1.0)\u00b2\nw = 0.1", "0.1", ACCENT_ORANGE,
     "Keeps noise magnitude near\ntarget (prevents scale drift)"),
    ("L_\u03ba", "Kappa Reg", "(\u03ba_mean - 2.0)\u00b2 + entropy\nw_pull=0.05, w_ent=0.01", "0.05", ACCENT_RED,
     "Pull \u03ba above 1.0 (KL activates)\n+ encourage per-token diversity"),
    ("L_\u03c3", "Sigma Pull", "(\u03c3_mean - 0.3)\u00b2\nw = 0.1", "0.1", ACCENT_PURPLE,
     "Prevent \u03c3 collapse to floor\n(lower \u03c3 = lower MSE = trap)"),
    ("L_div", "Diversity", "-std(\u03bc) tokens/frames\nw = 0.1", "0.1", LIGHT_GRAY,
     "Spatial + temporal diversity\nof adapter noise patterns"),
]

for i, (name, title, formula, weight, clr, desc) in enumerate(loss_cards):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.2 + row * 2.5)

    card = add_shape(slide, x, y, Inches(3.7), Inches(2.2),
                     fill_color=BG_CARD, line_color=clr, line_width=Pt(2))

    add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.4), Inches(0.35),
             f"{name}  \u2014  {title}", font_size=16, color=clr, bold=True)

    add_text(slide, x + Inches(0.15), y + Inches(0.5), Inches(3.4), Inches(0.7),
             formula, font_size=12, color=WHITE, font_name="Consolas")

    add_text(slide, x + Inches(0.15), y + Inches(1.3), Inches(3.4), Inches(0.7),
             desc, font_size=12, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: TRAINING DIAGNOSTICS & FIXES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.6),
         "Training Diagnostics: 3 Failure Modes & Fixes", font_size=36, color=WHITE, bold=True)

issues = [
    ("KL Permanently 0", ACCENT_RED,
     "Old formula: (D-1)/2 \u00b7 [log(\u03ba) - log(1+\u03ba)]",
     "Always negative for our parametrization\n(higher \u03ba = MORE spread, opposite of vMF)",
     "Fixed formula: (D-1)/2 \u00b7 log(1 + 1/\u03ba)\nAlways \u2265 0, correct sign convention",
     ACCENT_GREEN),
    ("Sigma Collapse to \u03c3_min", ACCENT_ORANGE,
     "sigma_mean: 0.495 \u2192 0.051 (~\u03c3_min=0.05)",
     "Lower \u03c3 = less noise = lower MSE\nSigma head finds trivial solution \u03c3\u21920",
     "Added sigma_mean_pull: (\u03c3_mean - 0.3)\u00b2 \u00d7 0.1\nCounteracts MSE incentive to minimize \u03c3",
     ACCENT_GREEN),
    ("\u2016\u03bc\u2016 Drift Above Target", ACCENT_PURPLE,
     "mu_norm_mean: 1.53 vs target 1.0",
     "magnitude_reg_weight=0.01 too weak\nvs flow matching loss ~0.016",
     "Increased magnitude_reg: 0.01 \u2192 0.1\n10x stronger penalty on scale drift",
     ACCENT_GREEN),
]

for i, (title, title_clr, symptom, cause, fix, fix_clr) in enumerate(issues):
    y = Inches(1.2 + i * 2.0)

    # Issue number
    num_circle = add_shape(slide, Inches(0.8), y, Inches(0.5), Inches(0.5),
                           fill_color=title_clr, shape_type=MSO_SHAPE.OVAL)
    p = num_circle.text_frame.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text(slide, Inches(1.5), y, Inches(3.0), Inches(0.5),
             title, font_size=20, color=title_clr, bold=True)

    # Symptom
    add_text(slide, Inches(1.5), y + Inches(0.45), Inches(3.0), Inches(0.4),
             symptom, font_size=11, color=GRAY, font_name="Consolas")

    # Cause box
    cause_box = add_shape(slide, Inches(4.8), y, Inches(3.5), Inches(1.0),
                          fill_color=BG_CARD, line_color=title_clr, line_width=Pt(1))
    add_text(slide, Inches(4.9), y - Inches(0.05), Inches(3.3), Inches(0.25),
             "Root Cause:", font_size=11, color=title_clr, bold=True)
    add_text(slide, Inches(4.9), y + Inches(0.25), Inches(3.3), Inches(0.7),
             cause, font_size=11, color=GRAY)

    # Arrow
    draw_right_arrow_shape(slide, Inches(8.3), y + Inches(0.5), Inches(0.4), color=ACCENT_GREEN)

    # Fix box
    fix_box = add_shape(slide, Inches(8.8), y, Inches(3.8), Inches(1.0),
                        fill_color=BG_CARD, line_color=fix_clr, line_width=Pt(1))
    add_text(slide, Inches(8.9), y - Inches(0.05), Inches(3.5), Inches(0.25),
             "Fix:", font_size=11, color=fix_clr, bold=True)
    add_text(slide, Inches(8.9), y + Inches(0.25), Inches(3.5), Inches(0.7),
             fix, font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 11: GEODESIC DIVERSITY
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "Geodesic Diversity on S\u00b9\u00b2\u2077", font_size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "Angular distance between noise directions is geometrically meaningful \u2014 unlike Euclidean distance in R\u00b9\u00b2\u2078",
         font_size=16, color=GRAY)

# Sphere illustration (simplified)
sphere_cx = Inches(3.0)
sphere_cy = Inches(4.0)
sphere_r = Inches(1.8)

# Sphere outline
sphere = add_shape(slide, sphere_cx - sphere_r, sphere_cy - sphere_r,
                   sphere_r * 2, sphere_r * 2,
                   line_color=RGBColor(0x47, 0x55, 0x69), line_width=Pt(2),
                   shape_type=MSO_SHAPE.OVAL)

# Equator line (horizontal ellipse)
equator = add_shape(slide, sphere_cx - sphere_r, sphere_cy - Inches(0.2),
                    sphere_r * 2, Inches(0.4),
                    line_color=RGBColor(0x33, 0x44, 0x55), line_width=Pt(1),
                    shape_type=MSO_SHAPE.OVAL)

# Points on sphere (adapter directions)
points = [
    (Inches(2.3), Inches(3.2), "\u03bc\u0302\u2081", ACCENT_BLUE),
    (Inches(3.5), Inches(2.5), "\u03bc\u0302\u2082", ACCENT_GREEN),
    (Inches(4.0), Inches(3.8), "\u03bc\u0302\u2083", ACCENT_PURPLE),
    (Inches(2.8), Inches(4.8), "\u03bc\u0302\u2084", ACCENT_ORANGE),
]

for px, py, label, clr in points:
    dot = add_shape(slide, px, py, Inches(0.2), Inches(0.2),
                    fill_color=clr, shape_type=MSO_SHAPE.OVAL)
    add_text(slide, px + Inches(0.15), py - Inches(0.15), Inches(0.5), Inches(0.3),
             label, font_size=12, color=clr, bold=True, font_name="Consolas")

add_text(slide, Inches(1.2), Inches(5.5), Inches(4.0), Inches(0.8),
         "S\u00b9\u00b2\u2077 hypersphere\n(shown as S\u00b2 for visualization)",
         font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Right side: metrics
metrics_x = Inches(6.5)
add_text(slide, metrics_x, Inches(2.0), Inches(6.0), Inches(0.4),
         "Diversity Metrics", font_size=22, color=ACCENT_BLUE, bold=True)

metrics = [
    ("Geodesic Distance", "d(\u03bc\u0302_i, \u03bc\u0302_j) = arccos(\u03bc\u0302_i \u00b7 \u03bc\u0302_j)", "Angular distance in [0, \u03c0] radians"),
    ("Geodesic Diversity", "mean(d(\u03bc\u0302_i, \u03bc\u0302_j)) for all pairs", "Higher = more diverse directions"),
    ("Observed Value", "1.35 rad (\u224877\u00b0)", "Well-spread on S\u00b9\u00b2\u2077 (max = \u03c0/2 \u2248 90\u00b0 for random)"),
]

for i, (name, formula, desc) in enumerate(metrics):
    y = Inches(2.6 + i * 1.4)
    m_box = add_shape(slide, metrics_x, y, Inches(5.8), Inches(1.1),
                      fill_color=BG_CARD, line_color=ACCENT_BLUE, line_width=Pt(1))
    add_text(slide, metrics_x + Inches(0.2), y + Inches(0.05), Inches(5.4), Inches(0.3),
             name, font_size=16, color=ACCENT_BLUE, bold=True)
    add_text(slide, metrics_x + Inches(0.2), y + Inches(0.35), Inches(5.4), Inches(0.3),
             formula, font_size=14, color=WHITE, font_name="Consolas")
    add_text(slide, metrics_x + Inches(0.2), y + Inches(0.7), Inches(5.4), Inches(0.3),
             desc, font_size=13, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 12: 4-FRAME GENERATION PIPELINE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.6),
         "How v1f Generates 4 Frames in 1 Shot", font_size=36, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         "The adapter produces 1344 tokens of structured noise \u2014 336 tokens per frame, each with unique direction and concentration",
         font_size=16, color=GRAY)

# Frame-by-frame breakdown
for f_idx in range(4):
    x = Inches(0.5 + f_idx * 3.1)

    # Frame header
    add_text(slide, x, Inches(1.7), Inches(2.8), Inches(0.4),
             f"Frame {f_idx}  (tokens {f_idx*336}\u2013{(f_idx+1)*336-1})", font_size=14, color=ACCENT_BLUE, bold=True)

    # Noise direction box
    draw_box_with_label(slide, x, Inches(2.2), Inches(2.8), Inches(0.7),
                        f"\u03bc\u0302 directions\n336 unique dirs on S\u00b9\u00b2\u2077",
                        fill=BG_CARD, border=ACCENT_BLUE, font_size=11)

    draw_down_arrow(slide, x + Inches(1.4), Inches(2.9), Inches(0.2), color=ACCENT_GREEN)

    # Concentration box
    draw_box_with_label(slide, x, Inches(3.3), Inches(2.8), Inches(0.7),
                        f"\u03ba per token\nEdges: \u03ba=5\u201310, Flat: \u03ba=0.5\u20132",
                        fill=BG_CARD, border=ACCENT_ORANGE, font_size=11)

    draw_down_arrow(slide, x + Inches(1.4), Inches(4.0), Inches(0.2), color=ACCENT_GREEN)

    # Sigma box
    draw_box_with_label(slide, x, Inches(4.4), Inches(2.8), Inches(0.7),
                        f"\u03c3_i per token\n\u03c3 \u2208 [0.05, 0.95]",
                        fill=BG_CARD, border=ACCENT_PURPLE, font_size=11)

    draw_down_arrow(slide, x + Inches(1.4), Inches(5.1), Inches(0.2), color=ACCENT_GREEN)

    # Structured noise
    draw_box_with_label(slide, x, Inches(5.5), Inches(2.8), Inches(0.7),
                        f"z_i = r_i \u00b7 SpCauchy(\u03bc\u0302_i, \u03ba_i)\n336 structured noise tokens",
                        fill=BG_CARD, border=ACCENT_GREEN, font_size=11)

# Bottom: single DiT pass
all_frames_arrow_y = Inches(6.4)
for f_idx in range(4):
    x = Inches(0.5 + f_idx * 3.1)
    draw_down_arrow(slide, x + Inches(1.4), Inches(6.2), Inches(0.15), color=WHITE)

dit_box = add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
                    fill_color=BG_CARD, line_color=WHITE, line_width=Pt(2))
p = dit_box.text_frame.paragraphs[0]
p.text = "48-Layer DiT  \u2192  Single Forward Pass  \u2192  4 Clean Frames (768\u00d7448 each)"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════
# SLIDE 13: WANDB METRICS TO WATCH
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.6),
         "Key Metrics to Monitor (W&B)", font_size=36, color=WHITE, bold=True)

metrics = [
    ("vfm/loss_mf", "Flow matching MSE", "Should decrease steadily", "0.21 \u2192 0.016", "\u2193", ACCENT_BLUE),
    ("vfm/loss_kl", "Spherical KL", "Should be > 0 (was stuck at 0!)", "0 \u2192 8\u201380", "\u2191", ACCENT_GREEN),
    ("vfm/kappa_mean", "Mean concentration", "Should converge near target (2.0)", "0.1 \u2192 2.0+", "\u2192 2.0", ACCENT_ORANGE),
    ("vfm/kappa_std", "Concentration diversity", "Higher = more per-token variation", "0 \u2192 3.0+", "\u2191", ACCENT_ORANGE),
    ("vfm/sigma_mean", "Mean noise level", "Should stay near 0.3 (not 0.05!)", "0.5 \u2192 0.3", "\u2192 0.3", ACCENT_PURPLE),
    ("vfm/mu_norm_mean", "Noise magnitude", "Should converge to target (1.0)", "var \u2192 1.0", "\u2192 1.0", ACCENT_RED),
    ("vfm/geodesic_diversity", "Angular spread (rad)", "Higher = more diverse directions", "\u2192 1.35", "\u2191", LIGHT_GRAY),
]

for i, (metric, name, expected, observed, trend, clr) in enumerate(metrics):
    y = Inches(1.2 + i * 0.85)

    # Metric name
    add_text(slide, Inches(0.8), y, Inches(3.2), Inches(0.3),
             metric, font_size=14, color=clr, bold=True, font_name="Consolas")

    # Description
    add_text(slide, Inches(4.0), y, Inches(2.5), Inches(0.3),
             name, font_size=14, color=LIGHT_GRAY)

    # Expected behavior
    add_text(slide, Inches(6.5), y, Inches(3.5), Inches(0.3),
             expected, font_size=13, color=GRAY)

    # Observed
    add_text(slide, Inches(10.0), y, Inches(1.5), Inches(0.3),
             observed, font_size=13, color=clr, font_name="Consolas")

    # Trend arrow
    add_text(slide, Inches(11.8), y, Inches(0.7), Inches(0.3),
             trend, font_size=16, color=clr, bold=True, alignment=PP_ALIGN.CENTER)

# W&B link
add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
         "Live dashboard: https://wandb.ai/snoozie/vfm-v1f", font_size=14, color=ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# SLIDE 14: BENCHMARK RESULTS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "Benchmark: VFM vs LTX Desktop 8-Step", font_size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         "RTX 5090, image-to-video, excluding text encoding time", font_size=16, color=GRAY)

# Table header
headers = ["Config", "Desktop 8-step", "VFM 1-step", "Speedup"]
header_widths = [Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)]

for j, (h, w) in enumerate(zip(headers, header_widths)):
    x = Inches(1.0 + j * 2.7)
    hdr = add_shape(slide, x, Inches(2.0), Inches(2.5), Inches(0.6),
                    fill_color=ACCENT_BLUE)
    p = hdr.text_frame.paragraphs[0]
    p.text = h
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Table rows
rows = [
    ("5s 540p", "~33s", "~3.8s", "8.7x"),
    ("5s 720p", "~42s", "~4.5s", "9.3x"),
    ("5s 1080p", "~76s", "~6.0s", "12.7x"),
    ("10s 540p", "~44s", "~4.5s", "9.8x"),
    ("20s 540p", "~73s", "~3.4s", "21.5x"),
]

for i, (config, desktop, vfm, speedup) in enumerate(rows):
    y = Inches(2.65 + i * 0.55)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x1A, 0x25, 0x40)

    vals = [config, desktop, vfm, speedup]
    clrs = [LIGHT_GRAY, ACCENT_RED, ACCENT_GREEN, ACCENT_BLUE]

    for j, (val, clr) in enumerate(zip(vals, clrs)):
        x = Inches(1.0 + j * 2.7)
        cell = add_shape(slide, x, y, Inches(2.5), Inches(0.5), fill_color=bg)
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(15)
        p.font.color.rgb = clr
        p.font.bold = (j == 3)
        p.alignment = PP_ALIGN.CENTER

# Note about theoretical
note_box = add_shape(slide, Inches(1.0), Inches(5.6), Inches(10.8), Inches(1.2),
                     fill_color=RGBColor(0x1A, 0x2A, 0x4A), line_color=ACCENT_ORANGE, line_width=Pt(1))
add_multiline(slide, Inches(1.3), Inches(5.7), Inches(10.2), Inches(1.0), [
    ("Note: VFM times are theoretical projections", ACCENT_ORANGE, True, 16),
    ("8x fewer transformer forward passes. Actual speedup depends on adapter overhead (~1-5%)", GRAY, False, 14),
    ("and VAE decode time (constant). Longer videos benefit more (adapter is amortized).", GRAY, False, 14),
])


# ════════════════════════════════════════════════════════════════
# SLIDE 15: WHAT'S NEXT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "What's Next", font_size=40, color=WHITE, bold=True)

next_items = [
    ("Validate v1f Fixes", "Run 6d3uvx1g with corrected KL, sigma pull, magnitude reg",
     "Active", ACCENT_GREEN),
    ("Enable Trajectory Distillation", "Switch distill_mode: output_match once v1f metrics are healthy",
     "Next", ACCENT_BLUE),
    ("Benchmark Real Quality", "Run vfm_benchmark.py --spherical with trained checkpoint",
     "Blocked", ACCENT_ORANGE),
    ("SLERP Interpolation", "Experimental: use geodesic interpolation for flow matching x_t",
     "Research", ACCENT_PURPLE),
    ("Adversarial Post-Training", "Add discriminator for sharper 1-step outputs (DAPT approach)",
     "Future", GRAY),
    ("Combine v1e + v1f", "Content-adaptive router WITH spherical noise \u2014 best of both",
     "Future", GRAY),
]

for i, (title, desc, status, clr) in enumerate(next_items):
    y = Inches(1.4 + i * 0.95)

    # Status badge
    badge = add_shape(slide, Inches(0.8), y, Inches(1.2), Inches(0.45),
                      fill_color=clr if status in ["Active"] else BG_CARD,
                      line_color=clr, line_width=Pt(1))
    p = badge.text_frame.paragraphs[0]
    p.text = status
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE if status == "Active" else clr
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text(slide, Inches(2.2), y - Inches(0.05), Inches(4.0), Inches(0.35),
             title, font_size=18, color=WHITE, bold=True)

    # Description
    add_text(slide, Inches(2.2), y + Inches(0.3), Inches(9.5), Inches(0.35),
             desc, font_size=14, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_path = "/home/johndpope/Documents/GitHub/ltx2-castlehill/packages/ltx-trainer/VFM_v1f_Spherical_Cauchy.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
